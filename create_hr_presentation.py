"""
Generate AI Workforce Guardian PowerPoint presentation for HR.
Run: pip install python-pptx
Then: python create_hr_presentation.py
Output: AI_Workforce_Guardian_HR_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt

# Layout indices: 0=Title, 1=Title+Content, 2=Section Header, 3=Two Content, 4=Comparison, 5=Title Only, 6=Blank
def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
    return slide

def add_content_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for b in bullets:
        p = body.add_paragraph()
        p.text = b
        p.level = 0
        p.space_after = Pt(6)
    return slide

def add_table_slide(prs, title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide.shapes.title.text = title
    # Add table
    left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(1) + Inches(0.3 * (len(rows) + 1))
    table = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height).table
    for i, h in enumerate(headers):
        table.cell(0, i).text = h
        table.cell(0, i).text_frame.paragraphs[0].font.bold = True
    for r_idx, row in enumerate(rows):
        for c_idx, cell_val in enumerate(row):
            table.cell(r_idx + 1, c_idx).text = str(cell_val)
    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs, "AI Workforce Guardian", "Employee Attrition Prediction & Analytics\nHR Analytics Dashboard")

    # Slide 2: Agenda
    add_content_slide(prs, "Today's Agenda", [
        "The problem: Why attrition matters",
        "Our solution: AI Workforce Guardian",
        "Key features & capabilities",
        "Dashboard walkthrough",
        "Model accuracy & reliability",
        "How HR can use it",
        "Next steps"
    ])

    # Slide 3: The Problem
    add_content_slide(prs, "The Cost of Employee Attrition", [
        "Financial Impact: Recruitment costs 50–75% of annual salary per hire; 6–9 months to full productivity",
        "Operational Impact: Project delays, team morale decline, institutional knowledge loss",
        "Industry average attrition rate: 13–15% annually",
        "Key Message: Acting after employees leave is expensive—we need to predict and prevent"
    ])

    # Slide 4: Reactive vs Proactive
    add_table_slide(prs, "From Reactive to Proactive Retention", 
        ["Approach", "Traditional (Reactive)", "AI-Powered (Proactive)"],
        [
            ("When", "Act after resignation", "Identify at-risk employees early"),
            ("Intervention", "Exit interviews (too late)", "Early retention conversations"),
            ("Programs", "Generic retention", "Targeted, personalized actions"),
            ("Insights", "Limited data", "Data-driven predictions"),
        ]
    )

    # Slide 5: Our Solution
    add_content_slide(prs, "AI Workforce Guardian — End-to-End Platform", [
        "Company-level dashboard: KPIs, attrition rates, department trends",
        "Individual risk analyzer: Per-employee attrition probability",
        "AI explanations: Why an employee is at risk",
        "HR recommendations: Actionable retention suggestions",
        "Employees at risk: Prioritized list for intervention"
    ])

    # Slide 6: Technology Stack
    add_content_slide(prs, "Built With Industry-Standard Tools", [
        "Machine Learning: Random Forest & SVM models; Scikit-learn; ~85% accuracy",
        "Analytics: Streamlit dashboard, Plotly charts, dark Netflix-style UI",
        "Data: IBM HR Analytics dataset (1,470 employees, 35+ features)",
        "Features: Age, Department, Job Role, Salary, Overtime, Satisfaction, Tenure, etc."
    ])

    # Slide 7: Key Metrics
    add_table_slide(prs, "Company-Level Metrics at a Glance",
        ["Metric", "What It Tells HR"],
        [
            ("Total Employees", "Current workforce size"),
            ("Attrition Rate %", "% who left in the period"),
            ("Employees at Risk", "Count with ≥30% predicted attrition probability"),
            ("Avg Monthly Income", "Compensation benchmark"),
        ]
    )

    # Slide 8: Risk Levels
    add_table_slide(prs, "How We Define Risk",
        ["Risk Level", "Probability", "HR Action"],
        [
            ("Low", "< 30%", "Monitor routinely"),
            ("Medium", "30–60%", "Schedule 1:1, review needs"),
            ("High", "≥ 60%", "Prioritize retention, immediate intervention"),
        ]
    )

    # Slide 9: AI Explanations
    add_content_slide(prs, "AI Explanations & HR Recommendations", [
        "AI Explanations (examples): 'High risk due to overtime, low salary'; 'Elevated risk from low job satisfaction'",
        "HR Recommendations: Reduce overtime; Review salary; Offer promotions; Improve work-life balance",
        "Not just a number—understand the why behind each prediction"
    ])

    # Slide 10: Top Factors
    add_content_slide(prs, "Top Factors Causing Attrition (From Model)", [
        "1. Overtime | 2. Monthly Income | 3. Job Satisfaction | 4. Years at Company",
        "5. Work-Life Balance | 6. Age | 7. Years in Current Role | 8. Environment Satisfaction",
        "Implication: Focus retention efforts on these levers"
    ])

    # Slide 11: Dashboard Sections
    add_table_slide(prs, "Dashboard Sections",
        ["Section", "Purpose"],
        [
            ("Company Dashboard", "KPIs, charts, overall trends"),
            ("HR Insights", "Department filter, top risk employees"),
            ("Individual Risk Analyzer", "Search by name, get prediction + explanation"),
            ("Attrition Data", "List of employees who left"),
            ("Employees at Risk", "Prioritized list; click for details"),
        ]
    )

    # Slide 12: Model Accuracy
    add_table_slide(prs, "Model Accuracy",
        ["Model", "Accuracy", "Notes"],
        [
            ("Random Forest", "~84%", "Primary model; feature importance"),
            ("SVM", "~85%", "Comparable performance"),
        ]
    )

    add_content_slide(prs, "Model Accuracy (continued)", [
        "Data split: 80% training, 20% testing (stratified for class balance)",
        "Predictions are probabilistic—use as a guide, not a guarantee"
    ])

    # Slide 14: How HR Can Use
    add_content_slide(prs, "Practical Use Cases for HR", [
        "Weekly review: Check Employees at Risk list; prioritize high-risk employees",
        "1:1 planning: Use AI recommendations to prepare talking points",
        "Department focus: Filter by department to see which teams need attention",
        "Retention programs: Target overtime, salary, satisfaction based on model insights",
        "Budget planning: Use risk counts for retention budget justification"
    ])

    # Slide 15: Data Privacy
    add_content_slide(prs, "Data Privacy & Ethical Use", [
        "Analysis uses anonymized/aggregated HR data",
        "Individual predictions for internal retention planning only",
        "No automated employment decisions—HR makes final calls",
        "Comply with local labor and data protection regulations"
    ])

    # Slide 16: Next Steps
    add_content_slide(prs, "Recommended Next Steps", [
        "1. Pilot: Run for 1–2 departments; validate with HR feedback",
        "2. Calibrate: Adjust risk threshold (e.g., 30%) based on capacity",
        "3. Integrate: Link to HRIS or performance data if available",
        "4. Train: Brief HR team on dashboard usage",
        "5. Review: Monthly review of model accuracy and outcomes"
    ])

    # Slide 17: Summary
    add_content_slide(prs, "Key Takeaways", [
        "AI Workforce Guardian predicts attrition risk before employees leave",
        "Dashboard provides company-level and individual-level insights",
        "AI explanations and recommendations support retention actions",
        "Model accuracy ~84–85%; use as decision support",
        "HR can act proactively to retain talent"
    ])

    # Slide 18: Thank You
    add_title_slide(prs, "Thank You", "Questions & Discussion")

    # Save
    out_path = "AI_Workforce_Guardian_HR_Presentation.pptx"
    prs.save(out_path)
    print(f"Created: {out_path}")

if __name__ == "__main__":
    main()
